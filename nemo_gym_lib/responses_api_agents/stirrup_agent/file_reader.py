# Copyright (c) 2025, NVIDIA CORPORATION.  All rights reserved.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
"""Extract text content and convert deliverable files for reward computation.

Reads common document formats (.docx, .pdf, .xlsx, .pptx, .txt, etc.)
and returns their combined text so the LLM judge can score actual content
rather than just the agent's finish-tool summary.

Also provides PDF conversion via LibreOffice headless for visual judging
with multimodal models (e.g., Gemini 3 Pro).
"""

from __future__ import annotations

import base64
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any


MAX_TOTAL_CHARS = 20_000


def read_deliverable_files(output_dir: str) -> str:
    """Read text from all deliverable files in *output_dir*.

    Returns a single string with sections like::

        === report.docx ===
        <extracted text>

        === data.xlsx ===
        <extracted text>

    Truncated to ~20k chars to stay within judge context limits.
    """
    output_path = Path(output_dir)
    if not output_path.is_dir():
        return ""

    files = sorted(output_path.iterdir())
    if not files:
        return ""

    parts: list[str] = []
    total_len = 0

    for fpath in files:
        if not fpath.is_file():
            continue

        ext = fpath.suffix.lower()
        try:
            text = _extract_text(fpath, ext)
        except Exception as exc:
            text = f"[Error reading {fpath.name}: {exc}]"

        if not text:
            continue

        section = f"=== {fpath.name} ===\n{text}"
        if total_len + len(section) > MAX_TOTAL_CHARS:
            remaining = MAX_TOTAL_CHARS - total_len
            if remaining > 200:
                section = section[:remaining] + "\n[...truncated]"
                parts.append(section)
            break

        parts.append(section)
        total_len += len(section)

    return "\n\n".join(parts)


def _extract_text(fpath: Path, ext: str) -> str:
    """Dispatch to the right extractor based on file extension."""
    if ext in (".txt", ".md", ".csv", ".json", ".html", ".xml", ".log"):
        return _read_text(fpath)
    elif ext == ".docx":
        return _read_docx(fpath)
    elif ext == ".pdf":
        return _read_pdf(fpath)
    elif ext == ".xlsx":
        return _read_xlsx(fpath)
    elif ext == ".pptx":
        return _read_pptx(fpath)
    else:
        size = os.path.getsize(fpath)
        return f"[Binary file: {fpath.name}, {size} bytes]"


def _read_text(fpath: Path) -> str:
    return fpath.read_text(encoding="utf-8", errors="replace").strip()


def _read_docx(fpath: Path) -> str:
    from docx import Document

    doc = Document(str(fpath))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_pdf(fpath: Path) -> str:
    from pdfminer.high_level import extract_text

    return extract_text(str(fpath)).strip()


def _read_xlsx(fpath: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(fpath), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(", ".join(cells))
        if rows:
            parts.append(f"Sheet: {sheet_name}\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def _read_pptx(fpath: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(fpath))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        if texts:
            parts.append(f"Slide {i}:\n" + "\n".join(texts))
    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# PDF conversion for visual judging (Gemini 3 Pro)
# ---------------------------------------------------------------------------

OFFICE_EXTS = {".docx", ".pptx", ".xlsx"}
TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".xml", ".html", ".yaml", ".yml", ".py", ".sh", ".log"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".heic", ".heif"}

MIME_TYPES = {
    ".pdf": "application/pdf",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
    ".heic": "image/heic",
    ".heif": "image/heif",
}


def _convert_office_to_pdf(fpath: Path) -> Path | None:
    """Convert a .docx/.xlsx/.pptx file to PDF using LibreOffice headless.

    Returns the path to the generated PDF, or None on failure.
    Uses a unique user profile to avoid lock conflicts in concurrent workers.

    Whitespace in the input filename makes LibreOffice's batch-convert mode
    silently drop the file (the URI it builds isn't percent-encoded), so we
    stage the input to a tempdir with a sanitized basename and move the PDF
    back to the original location.
    """
    profile_dir = Path(tempfile.mkdtemp(prefix="lo-profile-"))
    user_install = f"file://{profile_dir.as_posix()}"
    out_pdf = fpath.with_suffix(".pdf")
    stage_dir: Path | None = None
    input_path = fpath
    has_whitespace = any(c.isspace() for c in fpath.name)

    try:
        if has_whitespace:
            stage_dir = Path(tempfile.mkdtemp(prefix="lo-stage-"))
            safe_name = re.sub(r"\s+", "_", fpath.stem) + fpath.suffix
            input_path = stage_dir / safe_name
            shutil.copy2(fpath, input_path)
            lo_outdir = str(stage_dir)
        else:
            lo_outdir = str(fpath.parent)

        cmd = [
            "libreoffice",
            "--headless",
            "--nologo",
            "--nolockcheck",
            "--nodefault",
            "--norestore",
            f"-env:UserInstallation={user_install}",
            "--convert-to",
            "pdf",
            "--outdir",
            lo_outdir,
            str(input_path),
        ]
        p = subprocess.run(cmd, check=False, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=120)
        if stage_dir is not None:
            staged_pdf = stage_dir / (input_path.stem + ".pdf")
            if staged_pdf.exists():
                shutil.move(str(staged_pdf), str(out_pdf))
        if p.returncode != 0 or not out_pdf.exists():
            print(f"[file_reader] LibreOffice conversion failed for {fpath.name}: {p.stderr[:200]}", flush=True)
            return None
        return out_pdf
    except subprocess.TimeoutExpired:
        print(f"[file_reader] LibreOffice conversion timed out for {fpath.name}", flush=True)
        return None
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)
        if stage_dir is not None:
            shutil.rmtree(stage_dir, ignore_errors=True)


def convert_deliverables_to_content_blocks(output_dir: str) -> list[dict[str, Any]]:
    """Convert deliverable files to OpenAI-compatible content blocks for multimodal judging.

    Returns a list of content blocks suitable for the ``content`` field of an
    OpenAI chat message. Each file becomes either:

    - A text block (for .txt/.md/.csv etc.)
    - An image_url block with base64 data URL (for PDFs, images, converted Office docs)

    Office documents (.docx/.xlsx/.pptx) are first converted to PDF via
    LibreOffice headless so the judge sees rendered formatting/tables/charts.
    """
    output_path = Path(output_dir)
    if not output_path.is_dir():
        return []

    blocks: list[dict[str, Any]] = []
    converted_pdfs: list[Path] = []  # track for cleanup

    for fpath in sorted(output_path.iterdir()):
        if not fpath.is_file():
            continue

        ext = fpath.suffix.lower()

        try:
            if ext in TEXT_EXTS:
                text = fpath.read_text(encoding="utf-8", errors="replace").strip()
                if text:
                    blocks.append({"type": "text", "text": f"\n{fpath.name}:\n{text}"})

            elif ext in OFFICE_EXTS:
                pdf_path = _convert_office_to_pdf(fpath)
                if pdf_path and pdf_path.exists():
                    converted_pdfs.append(pdf_path)
                    data = pdf_path.read_bytes()
                    b64 = base64.b64encode(data).decode("ascii")
                    blocks.append({"type": "text", "text": f"\n{fpath.name} (converted to PDF):"})
                    blocks.append(
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:application/pdf;base64,{b64}"},
                        }
                    )
                else:
                    # Fallback to text extraction
                    text = _extract_text(fpath, ext)
                    if text:
                        blocks.append({"type": "text", "text": f"\n{fpath.name} (text fallback):\n{text}"})

            elif ext == ".pdf":
                data = fpath.read_bytes()
                b64 = base64.b64encode(data).decode("ascii")
                blocks.append({"type": "text", "text": f"\n{fpath.name}:"})
                blocks.append(
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:application/pdf;base64,{b64}"},
                    }
                )

            elif ext in IMAGE_EXTS:
                mime = MIME_TYPES.get(ext, "image/png")
                data = fpath.read_bytes()
                b64 = base64.b64encode(data).decode("ascii")
                blocks.append({"type": "text", "text": f"\n{fpath.name}:"})
                blocks.append(
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime};base64,{b64}"},
                    }
                )
        except Exception as exc:
            blocks.append({"type": "text", "text": f"\n{fpath.name}: [Error: {exc}]"})

    # Clean up converted PDFs (they live next to the originals)
    for pdf_path in converted_pdfs:
        pdf_path.unlink(missing_ok=True)

    return blocks
